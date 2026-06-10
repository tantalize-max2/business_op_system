# -*- coding: utf-8 -*-
"""
PPT通报生成模块 - 从gen_biz_report.py移植
支持：上传数据Excel + PPT模板 → 自动填充生成PPT
"""
import os
import re
import sys
import json
import tempfile
from datetime import datetime, timedelta
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from config import PPT_TEMPLATES_DIR, PPT_DATA_DIR


# ========== 模板管理 ==========

def ppt_template_path(name):
    safe_name = re.sub(r'[^\w\u4e00-\u9fff\-\.]', '_', name)
    return os.path.join(PPT_TEMPLATES_DIR, f"{safe_name}.json")


def list_ppt_templates():
    templates = []
    if os.path.exists(PPT_TEMPLATES_DIR):
        for fname in os.listdir(PPT_TEMPLATES_DIR):
            if not fname.endswith('.json'):
                continue
            fpath = os.path.join(PPT_TEMPLATES_DIR, fname)
            try:
                with open(fpath, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                templates.append({
                    'name': data.get('name', fname[:-5]),
                    'savedAt': data.get('savedAt', 0),
                    'hasDataFile': bool(data.get('dataFileData')),
                    'hasTemplate': bool(data.get('templateData'))
                })
            except Exception:
                pass
    templates.sort(key=lambda t: t.get('savedAt', 0), reverse=True)
    return templates


def save_ppt_template(name, template_data=None, data_file_data=None, data_map=None):
    template_record = {
        'name': name,
        'templateData': template_data or '',
        'dataFileData': data_file_data or '',
        'dataMap': data_map or {},
        'savedAt': datetime.now().timestamp() * 1000
    }
    fpath = ppt_template_path(name)
    with open(fpath, 'w', encoding='utf-8') as f:
        json.dump(template_record, f, ensure_ascii=False)
    return {'ok': True, 'name': name}


def get_ppt_template(name):
    fpath = ppt_template_path(name)
    if not os.path.exists(fpath):
        return None
    with open(fpath, 'r', encoding='utf-8') as f:
        return json.load(f)


def delete_ppt_template(name):
    fpath = ppt_template_path(name)
    if not os.path.exists(fpath):
        return False
    os.remove(fpath)
    return True


# ========== 上次标准化输出 ==========

def save_last_nz_output(output_path):
    """保存最近一次标准化填充的输出路径"""
    state_file = os.path.join(PPT_DATA_DIR, 'last_nz_output.json')
    with open(state_file, 'w', encoding='utf-8') as f:
        json.dump({'path': output_path, 'time': datetime.now().isoformat()}, f)


def get_last_nz_output():
    """获取最近一次标准化填充的输出路径"""
    state_file = os.path.join(PPT_DATA_DIR, 'last_nz_output.json')
    if not os.path.exists(state_file):
        return None
    try:
        with open(state_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
        path = data.get('path', '')
        if path and os.path.exists(path):
            return data
    except Exception:
        pass
    return None


# ========== 数据读取辅助 ==========

def cell_val(ws, col, row):
    v = ws.cell(row=row, column=col).value
    return v if v is not None else ''


def excel_date_to_str(serial):
    if isinstance(serial, (int, float)):
        base = datetime(1899, 12, 30)
        d = base + timedelta(days=int(serial))
        return f"{d.year}年{d.month}月{d.day}日"
    return str(serial)


def read_range(ws, min_row, max_row, min_col, max_col):
    return [[ws.cell(row=r, column=c).value for c in range(min_col, max_col+1)]
            for r in range(min_row, max_row+1)]


def fmt_pct(v):
    try:
        return f'{float(v)*100:.2f}%'
    except (ValueError, TypeError):
        return '0.00%'


def fmt_num(v, d=2):
    if v is None:
        return '0'
    if isinstance(v, str):
        return v
    try:
        fv = float(v)
        if d == 0 and fv == int(fv):
            return str(int(fv))
        return f'{fv:.{d}f}'
    except (ValueError, TypeError):
        return str(v) if v else ''


def fmt_comma(v):
    try:
        v = float(v)
        return f'{v:,.2f}' if v >= 10000 else f'{v:.2f}'
    except (ValueError, TypeError):
        return str(v)


def fmt_ca(s):
    """格式化数量(金额)字符串：如果数量和金额都为0，只输出'0'；去除内部多余空格"""
    if s is None:
        return ''
    s_str = str(s).strip().replace(' ', '')
    if not s_str:
        return ''
    c, a = parse_ca(s_str)
    if c == 0 and a == 0.0:
        return '0'
    return s_str


def _delivered_color_indices(data, col_idx):
    """根据指定列的数值，返回 {行索引: font_color} 标注前三深红、倒数两名绿"""
    DARK_RED = (0xC0, 0x00, 0x00)
    GREEN = (0x00, 0x80, 0x00)
    # 按 col_idx 数值排序，记录原始索引
    indexed = [(i, float(row[col_idx] or 0)) for i, row in enumerate(data) if row[col_idx] is not None]
    if not indexed:
        return {}
    indexed.sort(key=lambda x: x[1], reverse=True)
    result = {}
    # 前三标深红
    for i, _ in indexed[:min(3, len(indexed))]:
        result[i] = DARK_RED
    # 倒数两名标绿（取排序末尾，且不与标红重叠）
    bottom = indexed[-min(2, len(indexed)):]
    for i, _ in bottom:
        if i not in result:
            result[i] = GREEN
    return result


def get_bottom3(data):
    names = []
    for r in data[-3:]:
        m = re.match(r'(.+?)(（\d+）)?$', str(r[1]))
        names.append(m.group(1) if m else str(r[1]))
    return '、'.join(names)


def parse_ca(s):
    if s is None or str(s).strip() in ('0', ''):
        return (0, 0.0)
    s = str(s).strip().replace('，', ',').replace(' ', '')
    m = re.match(r'^(\d+)[\(（]([\d,.]+)[\)）]$', s)
    if m:
        return (int(m.group(1)), float(m.group(2).replace(',', '')))
    try:
        return (int(float(s)), 0.0)
    except (ValueError, TypeError):
        return (0, 0.0)


def calc_progress_summary(data):
    sq = sz = 0
    idc_c = idc_a = inet_c = inet_a = base_c = base_a = ict_c = ict_a = res_c = res_a = 0
    for row in data:
        try:
            sq += sum(int(float(row[i] or 0)) for i in [2, 3, 4])
            sz += int(float(row[5] or 0))
            c, a = parse_ca(row[6]); idc_c += c; idc_a += a
            c, a = parse_ca(row[7]); inet_c += c; inet_a += a
            c, a = parse_ca(row[8]); base_c += c; base_a += a
            c, a = parse_ca(row[9]); ict_c += c; ict_a += a
            c, a = parse_ca(row[10]); res_c += c; res_a += a
        except (ValueError, TypeError, IndexError):
            continue
    return dict(total=sq + sz, sq=sq, sz=sz,
                idc_c=idc_c, idc_a=idc_a, inet_c=inet_c, inet_a=inet_a,
                base_c=base_c, base_a=base_a, ict_c=ict_c, ict_a=ict_a,
                res_c=res_c, res_a=res_a)


def calc_effective_summary(eff_data):
    """计算有效商机汇总：分局金额之和 / 目标之和 → 完成率
    
    Returns:
        (sum_amount, sum_target, rate_pct, rate_text)
    """
    sum_amount = sum(float(item[1] or 0) for item in eff_data)
    sum_target = sum(float(item[2] or 0) for item in eff_data)
    bureau_count = len([item for item in eff_data if item[0] is not None and str(item[0]).strip()])
    if sum_target > 0:
        rate_pct = sum_amount / sum_target * 100
        rate_text = f'完成率{rate_pct:.2f}%（{fmt_num(sum_amount)}/{fmt_num(sum_target)}万）'
    else:
        rate_pct = 0
        rate_text = f'完成率0.00%（{fmt_num(sum_amount)}/0万）'
    return sum_amount, sum_target, rate_pct, rate_text, bureau_count


def count_below30(eff_data):
    """返回 (所有<30%分局数, 严重不足<10%分局名单)
    
    严重不足分局按完成率升序排列，最后一个名加"分局"
    """
    below30_count = 0
    severe = []  # (name, ratio)
    for item in eff_data:
        name, amt, tgt = str(item[0] or ''), float(item[1] or 0), float(item[2] or 0)
        if tgt <= 0:
            continue
        ratio = amt / tgt * 100
        if ratio < 30:
            below30_count += 1
        if ratio < 10:
            m = re.match(r'(.+?)(（\d+）)?$', name)
            clean_name = m.group(1) if m else name
            severe.append((clean_name, ratio))
    # 按完成率升序排列（最不足的排最前）
    severe.sort(key=lambda x: x[1])
    return below30_count, severe


def _format_severe_names(severe_list):
    """格式化严重不足分局名列表，最后一个加"分局"
    
    Args:
        severe_list: [(name, ratio), ...] 已排序
    Returns:
        str: "A、B、C分局" 格式
    """
    names = [item[0] for item in severe_list[:8]]
    if not names:
        return ''
    if len(names) == 1:
        return names[0] + '分局'
    return '、'.join(names[:-1]) + '、' + names[-1] + '分局'


# ========== PPT文本操作辅助 ==========

# 标题样式常量
TITLE_FONT = {'font_name': '微软雅黑', 'font_size': 28, 'font_bold': True}

def set_run_text(para, text):
    if not para.runs:
        return
    para.runs[0].text = str(text)
    for run in para.runs[1:]:
        run.text = ''


def set_shape_single_text(shape, text, font_name=None, font_size=None, font_bold=None, font_color=None):
    """设置shape中的文本，可选设置字体属性"""
    for para in shape.text_frame.paragraphs:
        if para.runs:
            set_run_text(para, '')
    if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
        run = shape.text_frame.paragraphs[0].runs[0]
        run.text = str(text)
        _apply_run_font(run, font_name, font_size, font_bold, font_color)
        # 确保所有段落的所有 run 都应用相同的加粗设置
        if font_bold is not None:
            for para in shape.text_frame.paragraphs:
                for r in para.runs:
                    r.font.bold = font_bold


def _apply_run_font(run, font_name=None, font_size=None, font_bold=None, font_color=None):
    """通用run字体设置"""
    if font_name:
        run.font.name = font_name
    if font_size is not None:
        run.font.size = Pt(font_size)
    if font_bold is not None:
        run.font.bold = font_bold
    if font_color:
        run.font.color.rgb = RGBColor(*font_color)


def set_shape_multiline(shape, lines, font_size=None, font_color_line1=None, font_color_line2=None):
    """设置多行文本，保留模板原始字号（font_size=None时）或统一设置字号
    
    font_color_line1: 第一行字体颜色 (R,G,B) 元组
    font_color_line2: 第二行字体颜色 (R,G,B) 元组
    """
    paras = shape.text_frame.paragraphs
    colors = [font_color_line1, font_color_line2]
    for i, line in enumerate(lines):
        line_color = colors[i] if i < len(colors) else None
        if i < len(paras):
            p = paras[i]
            if p.runs:
                r = p.runs[0]
                r.text = str(line)
                if font_size is not None:
                    r.font.size = Pt(font_size)
                if line_color:
                    r.font.color.rgb = RGBColor(*line_color)
                for run in p.runs[1:]:
                    run.text = ''
            else:
                r = p.add_run()
                r.text = str(line)
                if font_size is not None:
                    r.font.size = Pt(font_size)
                if line_color:
                    r.font.color.rgb = RGBColor(*line_color)
        else:
            p = shape.text_frame.add_paragraph()
            r = p.add_run()
            r.text = str(line)
            if font_size is not None:
                r.font.size = Pt(font_size)
            if line_color:
                r.font.color.rgb = RGBColor(*line_color)
    # 清除多余的空段落runs
    for i in range(len(lines), len(paras)):
        for run in paras[i].runs:
            run.text = ''


def _apply_cell_font(run, font_color=None, font_size=None):
    """统一设置单元格run字体：微软雅黑 加粗，可选颜色和字号覆盖"""
    run.font.name = '微软雅黑'
    run.font.size = Pt(font_size or 10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*font_color) if font_color else RGBColor(0x00, 0x00, 0x00)


def set_cell_text(cell, text, font_color=None, font_size=None):
    p = cell.text_frame.paragraphs[0] if cell.text_frame.paragraphs else None
    if p:
        if p.runs:
            run = p.runs[0]
            run.text = str(text)
            _apply_cell_font(run, font_color, font_size)
            for run in p.runs[1:]:
                run.text = ''
        else:
            run = p.add_run()
            run.text = str(text)
            _apply_cell_font(run, font_color, font_size)
    else:
        cell.text = str(text)
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            _apply_cell_font(cell.text_frame.paragraphs[0].runs[0], font_color, font_size)


def replace_picture(slide, shape_name, img_path):
    for shape in slide.shapes:
        if shape.name == shape_name:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(img_path, left, top, width, height)
            return True
    return False


def period_to_crm(period):
    m = re.match(r'(\d+)月(\d+)日[—\-](\d+)月(\d+)日', period)
    if m:
        return f'{m.group(1)}.{m.group(2)}-{m.group(3)}.{m.group(4)}'
    return period


def update_crm_date(shape, period_str, crm_date_full):
    """更新CRM日期：用正则匹配所有日期模式确保与标题时间一致
    
    period_str: 如 "1月1日—5月22日"
    crm_date_full: 如 "2026.1.1-2026.5.22"
    """
    for para in shape.text_frame.paragraphs:
        full_text = para.text
        if not full_text:
            continue
        new_text = full_text
        # 替换中文日期: X月X日—X月X日 或 X月X日-X月X日
        new_text = re.sub(r'\d+月\d+日[—\-]\d+月\d+日', period_str, new_text)
        # 替换完整CRM日期: XXXX.X.X-XXXX.X.XX (匹配任意年份)
        new_text = re.sub(r'\d{4}\.\d+\.\d+-\d{4}\.\d+\.\d+', crm_date_full, new_text)
        if new_text != full_text and para.runs:
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run.text = ''


# ========== 图表生成（调用 biz-chart 技能脚本） ==========

SCRIPTS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'scripts')
GEN_BIZ_CHART = os.path.join(SCRIPTS_DIR, 'gen_biz_chart.py')


def _generate_biz_chart(eff_data, output_path, title_prefix, date_str, year='2026'):
    """
    使用 biz-chart 技能的 gen_biz_chart.py 生成专业横向柱状图
    
    Args:
        eff_data: 有效商机数据列表 [[name, amount, target], ...]
        output_path: 输出PNG路径
        title_prefix: 标题前缀（"行业各分局" / "商业各分局"）
        date_str: 截止日期字符串
    """
    # 构造 gen_biz_chart.py 需要的 JSON 数据
    chart_items = []
    for item in eff_data:
        name = str(item[0] or '').strip()
        amt = float(item[1] or 0)
        tgt = float(item[2] or 0)
        # 第4列作为数量（如果有的话）
        cnt = 0
        if len(item) > 3 and item[3] is not None:
            try:
                cnt = int(float(item[3]))
            except (ValueError, TypeError):
                cnt = 0
        if not name or tgt <= 0:
            continue
        # 跳过合计行
        if name in ('合计', '总计', '小计', '行业合计', '商业合计'):
            continue
        chart_items.append({
            'name': name,
            'amount': round(amt, 2),
            'target': round(tgt, 2),
            'count': cnt
        })
    
    if not chart_items:
        return False
    
    # 写入临时 JSON
    json_path = output_path + '.json'
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(chart_items, f, ensure_ascii=False, indent=2)
    
    try:
        import subprocess
        cmd = [
            sys.executable, GEN_BIZ_CHART,
            '--data', json_path,
            '--output', output_path,
            '--date', date_str,
            '--title-prefix', title_prefix,
            '--title-suffix', '有效商机纳管情况',
            '--subtitle', '有效商机 / 商机储备目标',
            '--xlabel', '有效商机完成率',
            '--target-line', '30',
            '--year', year,
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=30, encoding='utf-8')
        if result.returncode != 0:
            print(f'Chart gen error: {result.stderr}')
            return False
        return os.path.exists(output_path)
    except Exception as e:
        print(f'Chart generation failed: {e}')
        return False
    finally:
        try:
            os.unlink(json_path)
        except Exception:
            pass


# ========== 核心：PPT生成 ==========

# ========== 默认数据映射（原始Excel行列位置） ==========
DEFAULT_DATA_MAP = {
    'date_cell': 'B30',          # 截止日期
    'period_cell': 'B31',        # 周期（如 1月1日—5月22日）
    'crm_date_cell': 'C30',      # 商机创建时间（如 2026.2.4-2026.6.25）
    'B27_cell': 'B27',           # 行业储备说明
    'B28_cell': 'B28',           # 商业储备说明
    'J27_cell': 'J27',           # 行业交付说明1
    'J28_cell': 'J28',           # 商业交付说明1
    'AI27_cell': 'AI27',         # 行业交付说明2
    'AI28_cell': 'AI28',         # 商业交付说明2
    'industry_reserve': 'A2:G12',      # 行业储备
    'commercial_reserve': 'A13:G25',    # 商业储备
    'industry_effective': 'J2:M13',     # 行业有效商机
    'commercial_effective': 'J14:M25',  # 商业有效商机
    'industry_progress': 'V5:AF15',     # 行业项目推进
    'commercial_progress': 'V16:AF28',  # 商业项目推进
    'industry_delivered': 'AH2:AK12',   # 行业交付
    'commercial_delivered': 'AH13:AK25', # 商业交付
}


def _parse_cell_ref(ref):
    """将单元格引用如 B30 转为 (col_1based, row_1based)"""
    m = re.match(r'^([A-Z]{1,3})(\d+)$', ref.upper().strip())
    if not m:
        return None, None
    col_str, row_str = m.group(1), m.group(2)
    col = 0
    for ch in col_str:
        col = col * 26 + (ord(ch) - ord('A') + 1)
    return col, int(row_str)


def _parse_range_ref(ref):
    """将范围引用如 A2:G12 转为 (min_row, max_row, min_col, max_col)，1-based"""
    m = re.match(r'^([A-Z]{1,3})(\d+):([A-Z]{1,3})(\d+)$', ref.upper().strip())
    if not m:
        return None, None, None, None
    c1, r1 = _parse_cell_ref(m.group(1) + m.group(2))
    c2, r2 = _parse_cell_ref(m.group(3) + m.group(4))
    if c1 is None or c2 is None:
        return None, None, None, None
    return min(r1, r2), max(r1, r2), min(c1, c2), max(c1, c2)


def generate_ppt(template_bytes, data_bytes, custom_texts=None, data_map=None):
    """
    根据模板PPT和数据Excel生成通报PPT

    Args:
        template_bytes: PPT模板文件字节
        data_bytes: 数据Excel文件字节
        custom_texts: 自定义文本替换（可选）
        data_map: 数据映射配置（可选），指定Excel中各区域位置

    Returns:
        dict: {ok: True, output_path: str} 或 {ok: False, error: str}
    """
    dm = dict(DEFAULT_DATA_MAP)
    if data_map:
        dm.update(data_map)

    # 解析单元格引用
    date_col, date_row = _parse_cell_ref(dm['date_cell'])
    period_col, period_row = _parse_cell_ref(dm['period_cell'])
    crm_col, crm_row = _parse_cell_ref(dm['crm_date_cell'])
    b27_col, b27_row = _parse_cell_ref(dm['B27_cell'])
    b28_col, b28_row = _parse_cell_ref(dm['B28_cell'])
    j27_col, j27_row = _parse_cell_ref(dm['J27_cell'])
    j28_col, j28_row = _parse_cell_ref(dm['J28_cell'])
    ai27_col, ai27_row = _parse_cell_ref(dm['AI27_cell'])
    ai28_col, ai28_row = _parse_cell_ref(dm['AI28_cell'])
    # 1. 写入临时文件
    tmp_ppt = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    tmp_ppt.write(template_bytes)
    tmp_ppt.close()

    tmp_data = tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False)
    tmp_data.write(data_bytes)
    tmp_data.close()

    tmp_out = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    tmp_out.close()

    # 2. 读取数据Excel
    try:
        wb = load_workbook(tmp_data.name, data_only=True)
        ws = wb.active
    except Exception as e:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except: pass
        return {'ok': False, 'error': f'读取数据文件失败: {str(e)}'}

    try:
        # 读取关键数据（使用可配置的单元格位置）
        DATE_STR = excel_date_to_str(cell_val(ws, date_col, date_row))
        PERIOD_STR = str(cell_val(ws, period_col, period_row))
        PERIOD_CRM = period_to_crm(PERIOD_STR)

        # 从日期(如"2026年5月22日")动态提取年份和短年份
        _year_match = re.match(r'(\d{4})年', DATE_STR)
        YEAR_FULL = _year_match.group(1) if _year_match else '2026'
        YEAR_SHORT = YEAR_FULL[2:]  # "26"

        # 商机创建时间：优先从C30读取，否则从周期派生
        CRM_DATE_C30 = str(cell_val(ws, crm_col, crm_row)).strip()
        if CRM_DATE_C30 and CRM_DATE_C30 != 'None':
            CRM_DATE_FULL = CRM_DATE_C30
        else:
            PERIOD_CRM = period_to_crm(PERIOD_STR)
            CRM_DATE_FULL = f'{YEAR_FULL}.{PERIOD_CRM}'.replace('-', f'-{YEAR_FULL}.')

        B27_TEXT = str(cell_val(ws, b27_col, b27_row)).strip().replace('\n', '')
        B28_TEXT = str(cell_val(ws, b28_col, b28_row)).strip().replace('\n', '')
        J27_TEXT = str(cell_val(ws, j27_col, j27_row)).strip()
        J28_TEXT = str(cell_val(ws, j28_col, j28_row)).strip()
        AI27_TEXT = str(cell_val(ws, ai27_col, ai27_row)).strip()
        AI28_TEXT = str(cell_val(ws, ai28_col, ai28_row)).strip()

        # 允许自定义文本覆盖
        if custom_texts:
            B27_TEXT = custom_texts.get('B27', B27_TEXT)
            B28_TEXT = custom_texts.get('B28', B28_TEXT)
            J27_TEXT = custom_texts.get('J27', J27_TEXT)
            J28_TEXT = custom_texts.get('J28', J28_TEXT)
            AI27_TEXT = custom_texts.get('AI27', AI27_TEXT)
            AI28_TEXT = custom_texts.get('AI28', AI28_TEXT)

        # 读取各区域数据（使用可配置的范围引用）
        ir_r1, ir_r2, ir_c1, ir_c2 = _parse_range_ref(dm['industry_reserve'])
        industry_reserve = read_range(ws, ir_r1, ir_r2, ir_c1, ir_c2)
        industry_reserve.sort(key=lambda x: float(x[6]) if x[6] is not None else 0, reverse=True)
        cr_r1, cr_r2, cr_c1, cr_c2 = _parse_range_ref(dm['commercial_reserve'])
        commercial_reserve = read_range(ws, cr_r1, cr_r2, cr_c1, cr_c2)
        commercial_reserve.sort(key=lambda x: float(x[6]) if x[6] is not None else 0, reverse=True)

        ie_r1, ie_r2, ie_c1, ie_c2 = _parse_range_ref(dm['industry_effective'])
        industry_effective = read_range(ws, ie_r1, ie_r2, ie_c1, ie_c2)
        ce_r1, ce_r2, ce_c1, ce_c2 = _parse_range_ref(dm['commercial_effective'])
        commercial_effective = read_range(ws, ce_r1, ce_r2, ce_c1, ce_c2)
        ip_r1, ip_r2, ip_c1, ip_c2 = _parse_range_ref(dm['industry_progress'])
        industry_progress = read_range(ws, ip_r1, ip_r2, ip_c1, ip_c2)
        cp_r1, cp_r2, cp_c1, cp_c2 = _parse_range_ref(dm['commercial_progress'])
        commercial_progress = read_range(ws, cp_r1, cp_r2, cp_c1, cp_c2)
        id_r1, id_r2, id_c1, id_c2 = _parse_range_ref(dm['industry_delivered'])
        industry_delivered = read_range(ws, id_r1, id_r2, id_c1, id_c2)
        cd_r1, cd_r2, cd_c1, cd_c2 = _parse_range_ref(dm['commercial_delivered'])
        commercial_delivered = read_range(ws, cd_r1, cd_r2, cd_c1, cd_c2)

        wb.close()
    except Exception as e:
        wb.close()
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except: pass
        return {'ok': False, 'error': f'解析数据失败: {str(e)}'}

    # 3. 生成图表（使用 biz-chart 技能脚本）
    chart_dir = os.path.join(PPT_DATA_DIR, 'charts')
    os.makedirs(chart_dir, exist_ok=True)
    industry_chart_path = os.path.join(chart_dir, 'industry_chart.png')
    commercial_chart_path = os.path.join(chart_dir, 'commercial_chart.png')

    # 清除旧的缓存图表文件，避免使用过期数据
    for old_chart in [industry_chart_path, commercial_chart_path]:
        if os.path.exists(old_chart):
            try: os.unlink(old_chart)
            except: pass

    # 从周期字符串提取日期，如"1月1日—5月22日" → "5月22日"
    date_suffix = PERIOD_STR.split('—')[-1] if '—' in PERIOD_STR else PERIOD_STR.split('-')[-1] if '-' in PERIOD_STR else PERIOD_STR

    # 行业有效商机图表（Slide 6）
    ind_chart_ok = _generate_biz_chart(industry_effective, industry_chart_path, '行业各分局', date_suffix, year=YEAR_FULL)
    # 商业有效商机图表（Slide 7）
    comm_chart_ok = _generate_biz_chart(commercial_effective, commercial_chart_path, '商业各分局', date_suffix, year=YEAR_FULL)

    # 调试信息：有效商机图表实际使用的数据
    def _eff_data_debug(eff_data, max_n=5):
        items = []
        for item in eff_data[:max_n]:
            name = str(item[0] or '').strip()
            amt = float(item[1] or 0)
            tgt = float(item[2] or 0)
            items.append(f'{name}: {amt}/{tgt}')
        return '; '.join(items) + (f' ...共{len(eff_data)}行' if len(eff_data) > max_n else '')

    chart_debug = {
        'industry_effective_range': dm['industry_effective'],
        'industry_data_preview': _eff_data_debug(industry_effective),
        'industry_chart_ok': ind_chart_ok,
        'commercial_effective_range': dm['commercial_effective'],
        'commercial_data_preview': _eff_data_debug(commercial_effective),
        'commercial_chart_ok': comm_chart_ok,
    }

    # 4. 加载PPT模板并填充
    try:
        prs = Presentation(tmp_ppt.name)
    except Exception as e:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except: pass
        return {'ok': False, 'error': f'读取PPT模板失败: {str(e)}'}

    # 4.1 模板结构验证——必须包含关键shape
    if len(prs.slides) < 15:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except: pass
        return {'ok': False, 'error': f'模板页面不足（{len(prs.slides)}页，需15页）。请使用正确的商机通报模板。'}

    required_shapes = {
        0: ['date'],            # Slide 1: 日期
        2: ['表格 14'],         # Slide 3: 行业储备表格
        5: ['Text 7', 'Text 49'],  # Slide 6: 行业有效商机
        6: ['Text 7', 'Text 49'],  # Slide 7: 商业有效商机
    }
    missing = []
    for slide_idx, names in required_shapes.items():
        existing = {s.name for s in prs.slides[slide_idx].shapes}
        for name in names:
            if name not in existing:
                missing.append(f'Slide{slide_idx+1}/{name}')
    if missing:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except: pass
        return {'ok': False, 'error': f'模板缺少关键元素：{", ".join(missing)}。请使用正确的商机通报模板(含表格、图表占位等)。'}

    # ====== Slide 1: 日期 ======
    for shape in prs.slides[0].shapes:
        if shape.name == 'date':
            set_shape_single_text(shape, DATE_STR, font_name='微软雅黑', font_size=18, font_bold=True)

    # ====== Slide 3: 行业储备 ======
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, f'1、行业板块---{YEAR_FULL}年商机储备情况（{PERIOD_STR}）', **TITLE_FONT)
        elif shape.name == '表格 14':
            tbl = shape.table
            for i, rd in enumerate(industry_reserve):
                r = i + 1
                set_cell_text(tbl.cell(r, 0), str(i + 1))
                for c in range(1, 7):
                    v = rd[c]
                    if c == 6:
                        set_cell_text(tbl.cell(r, c), fmt_pct(v))
                    elif c == 5:
                        set_cell_text(tbl.cell(r, c), fmt_num(v))
                    elif c in (3, 4):
                        set_cell_text(tbl.cell(r, c), fmt_num(v, 0))
                    else:
                        set_cell_text(tbl.cell(r, c), str(v) if v else '')
        elif shape.name == '文本框 4':
            b3 = get_bottom3(industry_reserve)
            set_shape_multiline(shape, [
                f'1、{B27_TEXT}',
                f'2、储备完成率不足的分局：{b3}'
            ], font_color_line1=None, font_color_line2=(0x00, 0x80, 0x00))
        elif shape.name == '文本框 5':
            update_crm_date(shape, PERIOD_STR, CRM_DATE_FULL)

    # ====== Slide 4: 商业储备 ======
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, f'1、商校板块---{YEAR_FULL}年商机储备情况（{PERIOD_STR}）', **TITLE_FONT)
        elif shape.name == '表格 14':
            tbl = shape.table
            for i, rd in enumerate(commercial_reserve):
                r = i + 1
                set_cell_text(tbl.cell(r, 0), str(i + 1))
                for c in range(1, 7):
                    v = rd[c]
                    if c == 6:
                        set_cell_text(tbl.cell(r, c), fmt_pct(v))
                    elif c == 5:
                        set_cell_text(tbl.cell(r, c), fmt_num(v))
                    elif c in (3, 4):
                        set_cell_text(tbl.cell(r, c), fmt_num(v, 0))
                    else:
                        set_cell_text(tbl.cell(r, c), str(v) if v else '')
        elif shape.name == '文本框 4':
            b3 = get_bottom3(commercial_reserve)
            set_shape_multiline(shape, [
                f'1、{B28_TEXT}',
                f'2、储备完成率不足的分局：{b3}'
            ], font_color_line1=None, font_color_line2=(0x00, 0x80, 0x00))
        elif shape.name == '文本框 5':
            update_crm_date(shape, PERIOD_STR, CRM_DATE_FULL)

    # ====== Slide 6: 行业有效商机 ======
    slide6 = prs.slides[5]
    ind_below30_count, ind_severe = count_below30(industry_effective)
    ind_sum_amt, ind_sum_tgt, ind_rate_pct, ind_rate_text, ind_bureau_count = calc_effective_summary(industry_effective)

    for shape in slide6.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, f'2、行业板块---有效商机纳管情况（{YEAR_SHORT}年{PERIOD_STR}）', **TITLE_FONT)
        elif shape.name == 'Text 7':
            set_shape_single_text(shape, fmt_num(ind_sum_amt) + '万')
        elif shape.name == 'Text 9':
            set_shape_multiline(shape, ['有效商机/商机储备目标', ind_rate_text])
        elif shape.name == 'Text 49':
            names_str = _format_severe_names(ind_severe)
            set_shape_multiline(shape, [
                f'{ind_below30_count}个分局未达30%储备率（共{ind_bureau_count}个）',
                f'其中{names_str}商机储备纳管严重不足'
            ], font_size=14)
        elif shape.name in ('文本框 123', '文本框 5'):
            update_crm_date(shape, PERIOD_STR, CRM_DATE_FULL)

    if os.path.exists(industry_chart_path):
        replace_picture(slide6, '图片 6', industry_chart_path)

    # ====== Slide 7: 商业有效商机 ======
    slide7 = prs.slides[6]
    comm_below30_count, comm_severe = count_below30(commercial_effective)
    comm_sum_amt, comm_sum_tgt, comm_rate_pct, comm_rate_text, comm_bureau_count = calc_effective_summary(commercial_effective)

    for shape in slide7.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, f'2、商业板块---有效商机纳管情况（{YEAR_SHORT}年{PERIOD_STR}）', **TITLE_FONT)
        elif shape.name == 'Text 7':
            set_shape_single_text(shape, fmt_num(comm_sum_amt) + '万')
        elif shape.name == 'Text 9':
            set_shape_multiline(shape, ['有效商机/商机储备目标', comm_rate_text])
        elif shape.name == 'Text 49':
            names_str = _format_severe_names(comm_severe)
            set_shape_multiline(shape, [
                f'{comm_below30_count}个分局均未达30%储备率（共{comm_bureau_count}个）',
                f'其中{names_str}商机纳管储备严重不足'
            ], font_size=14)
        elif shape.name in ('文本框 123', '文本框 5'):
            update_crm_date(shape, PERIOD_STR, CRM_DATE_FULL)

    if os.path.exists(commercial_chart_path):
        replace_picture(slide7, '图片 1', commercial_chart_path)

    # ====== Slide 10: 行业项目推进 ======
    slide10 = prs.slides[9]
    ind_s = calc_progress_summary(industry_progress)

    for shape in slide10.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, f'3、行业板块--商机项目推进情况（{YEAR_SHORT}年{PERIOD_STR}）', **TITLE_FONT)
        elif shape.name == '表格 27':
            tbl = shape.table
            for i, rd in enumerate(industry_progress):
                r = i + 2
                for c in range(11):
                    val = rd[c] if rd[c] is not None else ''
                    # 列6-10为商机类型数据，格式化 0(0.0) → 0
                    if c >= 6 and c <= 10:
                        set_cell_text(tbl.cell(r, c), fmt_ca(val))
                    else:
                        set_cell_text(tbl.cell(r, c), str(val))
        elif shape.name == '文本框 3':
            s = ind_s
            set_shape_multiline(shape, [
                f'行业板块在途有效商机{s["total"]}条，其中处于售前{s["sq"]}条、售中{s["sz"]}条。',
                f'行业板块在途有效商机类型分布情况：IDC、互联网专线商机数量为0',
                f'ICT商机数量为{s["ict_c"]}，金额{fmt_comma(s["ict_a"])}万元，基础业务数量为{s["base_c"]}，金额{fmt_comma(s["base_a"])}万元，产数资源商机数量为{s["res_c"]}，金额{fmt_comma(s["res_a"])}万元。'
            ])

    # ====== Slide 11: 商业项目推进 ======
    slide11 = prs.slides[10]
    comm_s = calc_progress_summary(commercial_progress)

    for shape in slide11.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, f'3、商校板块--商机项目推进情况（{YEAR_SHORT}年{PERIOD_STR}）', **TITLE_FONT)
        elif shape.name == '表格 1':
            tbl = shape.table
            for i, rd in enumerate(commercial_progress):
                r = i + 2
                for c in range(11):
                    val = rd[c] if rd[c] is not None else ''
                    # 列6-10为商机类型数据，格式化 0(0.0) → 0
                    if c >= 6 and c <= 10:
                        set_cell_text(tbl.cell(r, c), fmt_ca(val))
                    else:
                        set_cell_text(tbl.cell(r, c), str(val))
        elif shape.name == '文本框 5':
            s = comm_s
            set_shape_multiline(shape, [
                f'商校板块在途有效商机{s["total"]}条，其中处于售前{s["sq"]}条、售中{s["sz"]}条。',
                f'商业板块有效商机类型分布情况：IDC、互联网专线商机数量为0',
                f'ICT商机数量为{s["ict_c"]}，金额{fmt_comma(s["ict_a"])}万元，基础业务数量为{s["base_c"]}，金额{fmt_comma(s["base_a"])}万元，产数资源商机数量为{s["res_c"]}，金额{fmt_comma(s["res_a"])}万元'
            ])

    # ====== Slide 13: 行业交付 ======
    slide13 = prs.slides[12]
    ind_count_colors = _delivered_color_indices(industry_delivered, 2)   # 已交付数量列
    ind_amount_colors = _delivered_color_indices(industry_delivered, 3)  # 已交付金额列

    for shape in slide13.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, f'3、行业板块--{YEAR_SHORT}年已完成交付项目情况（{YEAR_SHORT}年{PERIOD_STR}）', **TITLE_FONT)
        elif shape.name == '表格 27':
            tbl = shape.table
            for i, rd in enumerate(industry_delivered):
                r = i + 1
                for c in range(4):
                    v = rd[c]
                    color = None
                    if c == 2:
                        color = ind_count_colors.get(i)
                    elif c == 3:
                        color = ind_amount_colors.get(i)
                    # 黑色11号，彩色14号
                    fs = 14 if color else 11
                    if c == 3:
                        set_cell_text(tbl.cell(r, c), fmt_num(v), font_color=color, font_size=fs)
                    elif c == 2:
                        set_cell_text(tbl.cell(r, c), fmt_num(v, 0), font_color=color, font_size=fs)
                    else:
                        set_cell_text(tbl.cell(r, c), str(v) if v else '', font_color=color, font_size=fs)
        elif shape.name == '文本框 3':
            set_shape_single_text(shape, AI27_TEXT)

    # ====== Slide 14: 商业交付 ======
    slide14 = prs.slides[13]
    comm_count_colors = _delivered_color_indices(commercial_delivered, 2)
    comm_amount_colors = _delivered_color_indices(commercial_delivered, 3)

    for shape in slide14.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, f'3、商校板块--{YEAR_SHORT}年已完成交付项目情况（{YEAR_SHORT}年{PERIOD_STR}）', **TITLE_FONT)
        elif shape.name == '表格 2':
            tbl = shape.table
            for i, rd in enumerate(commercial_delivered):
                r = i + 1
                for c in range(4):
                    v = rd[c]
                    color = None
                    if c == 2:
                        color = comm_count_colors.get(i)
                    elif c == 3:
                        color = comm_amount_colors.get(i)
                    # 黑色11号，彩色14号
                    fs = 14 if color else 11
                    if c == 3:
                        set_cell_text(tbl.cell(r, c), fmt_num(v), font_color=color, font_size=fs)
                    elif c == 2:
                        set_cell_text(tbl.cell(r, c), fmt_num(v, 0), font_color=color, font_size=fs)
                    else:
                        set_cell_text(tbl.cell(r, c), str(v) if v else '', font_color=color, font_size=fs)
        elif shape.name == '文本框 5':
            set_shape_single_text(shape, AI28_TEXT)

    # 5. 保存输出
    try:
        prs.save(tmp_out.name)
    except Exception as e:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except: pass
        return {'ok': False, 'error': f'保存PPT失败: {str(e)}'}

    # 清理临时输入文件（输出文件由调用方负责清理）
    try: os.unlink(tmp_ppt.name)
    except: pass
    try: os.unlink(tmp_data.name)
    except: pass

    return {
        'ok': True,
        'output_path': tmp_out.name,
        'date': DATE_STR,
        'period': PERIOD_STR,
        'industry_bureaus': len(industry_reserve),
        'commercial_bureaus': len(commercial_reserve),
        'industry_below30': ind_below30_count,
        'commercial_below30': comm_below30_count,
        'chart_debug': chart_debug,
    }
