# -*- coding: utf-8 -*-
"""PPT 通报生成服务层 - 从数据 Excel + PPT 模板生成通报 PPT 的全部业务逻辑

从原 models/ppt_model.py 迁移。本模块负责：
- 数据区域读取与格式化（储备/有效商机/项目推进/交付）
- 完成率/汇总统计计算
- PPT 文本/表格/图片填充
- 行业/商业有效商机图表生成（调用 scripts/gen_biz_chart.py）
- 数据预览（供 routes 调用校准数据匹配）

模板持久化（CRUD）与上次标准化输出状态存取由 models/ppt_model.py 负责。
"""
import os
import re
import sys
import json
import tempfile
from datetime import datetime, timedelta
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from config import PPT_DATA_DIR


# ========== 默认数据映射（原始Excel行列位置） ==========
DEFAULT_DATA_MAP = {
    'date_cell': 'B30',
    'period_cell': 'B31',
    'crm_date_cell': 'C30',
    'B27_cell': 'B27',
    'B28_cell': 'B28',
    'J27_cell': 'J27',
    'J28_cell': 'J28',
    'AI27_cell': 'AI27',
    'AI28_cell': 'AI28',
    'industry_reserve': 'A2:G12',
    'commercial_reserve': 'A13:G25',
    'industry_effective': 'J2:M13',
    'commercial_effective': 'J14:M25',
    'industry_progress': 'V5:AF15',
    'commercial_progress': 'V16:AF28',
    'industry_delivered': 'AH2:AK12',
    'commercial_delivered': 'AH13:AK25',
}


# ========== 数据读取辅助 ==========

def cell_val(ws, col, row):
    v = ws.cell(row=row, column=col).value
    return v if v is not None else ''


def excel_date_to_str(serial):
    if isinstance(serial, (int, float)):
        base = datetime(1899, 12, 30)
        d = base + timedelta(days=int(serial))
        return f"{d.year}年{d.month}月{d.day}日"
    return str(serial)


def read_range(ws, min_row, max_row, min_col, max_col):
    return [[ws.cell(row=r, column=c).value for c in range(min_col, max_col + 1)]
            for r in range(min_row, max_row + 1)]


def _parse_cell_ref(ref):
    m = re.match(r'^([A-Z]{1,3})(\d+)$', ref.upper().strip())
    if not m:
        return None, None
    col_str, row_str = m.group(1), m.group(2)
    col = 0
    for ch in col_str:
        col = col * 26 + (ord(ch) - ord('A') + 1)
    return col, int(row_str)


def _parse_range_ref(ref):
    m = re.match(r'^([A-Z]{1,3})(\d+):([A-Z]{1,3})(\d+)$', ref.upper().strip())
    if not m:
        return None, None, None, None
    c1, r1 = _parse_cell_ref(m.group(1) + m.group(2))
    c2, r2 = _parse_cell_ref(m.group(3) + m.group(4))
    if c1 is None or c2 is None:
        return None, None, None, None
    return min(r1, r2), max(r1, r2), min(c1, c2), max(c1, c2)


# ========== 格式化 ==========

def fmt_pct(v):
    try:
        return f'{float(v) * 100:.2f}%'
    except (ValueError, TypeError):
        return '0.00%'


def fmt_num(v, d=2):
    if v is None:
        return '0'
    if isinstance(v, str):
        return v
    try:
        fv = float(v)
        if d == 0 and fv == int(fv):
            return str(int(fv))
        return f'{fv:.{d}f}'
    except (ValueError, TypeError):
        return str(v) if v else ''


def fmt_comma(v):
    try:
        v = float(v)
        return f'{v:,.2f}' if v >= 10000 else f'{v:.2f}'
    except (ValueError, TypeError):
        return str(v)


def fmt_ca(s):
    if s is None:
        return ''
    s_str = str(s).strip().replace(' ', '')
    if not s_str:
        return ''
    c, a = parse_ca(s_str)
    if c == 0 and a == 0.0:
        return '0'
    return s_str


# ========== 统计计算 ==========

def parse_ca(s):
    if s is None or str(s).strip() in ('0', ''):
        return (0, 0.0)
    s = str(s).strip().replace('，', ',').replace(' ', '')
    m = re.match(r'^(\d+)[\(（]([\d,.]+)[\)）]$', s)
    if m:
        return (int(m.group(1)), float(m.group(2).replace(',', '')))
    try:
        return (int(float(s)), 0.0)
    except (ValueError, TypeError):
        return (0, 0.0)


def calc_progress_summary(data):
    sq = sz = 0
    idc_c = idc_a = inet_c = inet_a = base_c = base_a = ict_c = ict_a = res_c = res_a = 0
    for row in data:
        try:
            sq += sum(int(float(row[i] or 0)) for i in [2, 3, 4])
            sz += int(float(row[5] or 0))
            c, a = parse_ca(row[6]); idc_c += c; idc_a += a
            c, a = parse_ca(row[7]); inet_c += c; inet_a += a
            c, a = parse_ca(row[8]); base_c += c; base_a += a
            c, a = parse_ca(row[9]); ict_c += c; ict_a += a
            c, a = parse_ca(row[10]); res_c += c; res_a += a
        except (ValueError, TypeError, IndexError):
            continue
    return dict(total=sq + sz, sq=sq, sz=sz,
                idc_c=idc_c, idc_a=idc_a, inet_c=inet_c, inet_a=inet_a,
                base_c=base_c, base_a=base_a, ict_c=ict_c, ict_a=ict_a,
                res_c=res_c, res_a=res_a)


def calc_effective_summary(eff_data):
    sum_amount = sum(float(item[1] or 0) for item in eff_data)
    sum_target = sum(float(item[2] or 0) for item in eff_data)
    bureau_count = len([item for item in eff_data if item[0] is not None and str(item[0]).strip()])
    if sum_target > 0:
        rate_pct = sum_amount / sum_target * 100
        rate_text = f'完成率{rate_pct:.2f}%（{fmt_num(sum_amount)}/{fmt_num(sum_target)}万）'
    else:
        rate_pct = 0
        rate_text = f'完成率0.00%（{fmt_num(sum_amount)}/0万）'
    return sum_amount, sum_target, rate_pct, rate_text, bureau_count


def count_below30(eff_data):
    below30_count = 0
    severe = []
    for item in eff_data:
        name, amt, tgt = str(item[0] or ''), float(item[1] or 0), float(item[2] or 0)
        if tgt <= 0:
            continue
        ratio = amt / tgt * 100
        if ratio < 30:
            below30_count += 1
        if ratio < 10:
            m = re.match(r'(.+?)(（\d+）)?$', name)
            clean_name = m.group(1) if m else name
            severe.append((clean_name, ratio))
    severe.sort(key=lambda x: x[1])
    return below30_count, severe


def _format_severe_names(severe_list):
    names = [item[0] for item in severe_list[:8]]
    if not names:
        return ''
    if len(names) == 1:
        return names[0] + '分局'
    return '、'.join(names[:-1]) + '、' + names[-1] + '分局'


def _delivered_color_indices(data, col_idx):
    DARK_RED = (0xC0, 0x00, 0x00)
    GREEN = (0x00, 0x80, 0x00)
    indexed = [(i, float(row[col_idx] or 0)) for i, row in enumerate(data) if row[col_idx] is not None]
    if not indexed:
        return {}
    indexed.sort(key=lambda x: x[1], reverse=True)
    result = {}
    for i, _ in indexed[:min(3, len(indexed))]:
        result[i] = DARK_RED
    bottom = indexed[-min(2, len(indexed)):]
    for i, _ in bottom:
        if i not in result:
            result[i] = GREEN
    return result


def get_bottom3(data):
    names = []
    for r in data[-3:]:
        m = re.match(r'(.+?)(（\d+）)?$', str(r[1]))
        names.append(m.group(1) if m else str(r[1]))
    return '、'.join(names)


# ========== PPT文本操作辅助 ==========

TITLE_FONT = {'font_name': '微软雅黑', 'font_size': 28, 'font_bold': True}


def set_run_text(para, text):
    if not para.runs:
        return
    para.runs[0].text = str(text)
    for run in para.runs[1:]:
        run.text = ''


def _apply_run_font(run, font_name=None, font_size=None, font_bold=None, font_color=None):
    if font_name:
        run.font.name = font_name
    if font_size is not None:
        run.font.size = Pt(font_size)
    if font_bold is not None:
        run.font.bold = font_bold
    if font_color:
        run.font.color.rgb = RGBColor(*font_color)


def set_shape_single_text(shape, text, font_name=None, font_size=None, font_bold=None, font_color=None):
    for para in shape.text_frame.paragraphs:
        if para.runs:
            set_run_text(para, '')
    if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
        run = shape.text_frame.paragraphs[0].runs[0]
        run.text = str(text)
        _apply_run_font(run, font_name, font_size, font_bold, font_color)
        if font_bold is not None:
            for para in shape.text_frame.paragraphs:
                for r in para.runs:
                    r.font.bold = font_bold


def set_shape_multiline(shape, lines, font_size=None, font_color_line1=None, font_color_line2=None):
    paras = shape.text_frame.paragraphs
    colors = [font_color_line1, font_color_line2]
    for i, line in enumerate(lines):
        line_color = colors[i] if i < len(colors) else None
        if i < len(paras):
            p = paras[i]
            if p.runs:
                r = p.runs[0]
                r.text = str(line)
                if font_size is not None:
                    r.font.size = Pt(font_size)
                if line_color:
                    r.font.color.rgb = RGBColor(*line_color)
                for run in p.runs[1:]:
                    run.text = ''
            else:
                r = p.add_run()
                r.text = str(line)
                if font_size is not None:
                    r.font.size = Pt(font_size)
                if line_color:
                    r.font.color.rgb = RGBColor(*line_color)
        else:
            p = shape.text_frame.add_paragraph()
            r = p.add_run()
            r.text = str(line)
            if font_size is not None:
                r.font.size = Pt(font_size)
            if line_color:
                r.font.color.rgb = RGBColor(*line_color)
    for i in range(len(lines), len(paras)):
        for run in paras[i].runs:
            run.text = ''


def _apply_cell_font(run, font_color=None, font_size=None):
    run.font.name = '微软雅黑'
    run.font.size = Pt(font_size or 10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*font_color) if font_color else RGBColor(0x00, 0x00, 0x00)


def set_cell_text(cell, text, font_color=None, font_size=None):
    p = cell.text_frame.paragraphs[0] if cell.text_frame.paragraphs else None
    if p:
        if p.runs:
            run = p.runs[0]
            run.text = str(text)
            _apply_cell_font(run, font_color, font_size)
            for run in p.runs[1:]:
                run.text = ''
        else:
            run = p.add_run()
            run.text = str(text)
            _apply_cell_font(run, font_color, font_size)
    else:
        cell.text = str(text)
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            _apply_cell_font(cell.text_frame.paragraphs[0].runs[0], font_color, font_size)


def replace_picture(slide, shape_name, img_path):
    for shape in slide.shapes:
        if shape.name == shape_name:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(img_path, left, top, width, height)
            return True
    return False


def period_to_crm(period):
    m = re.match(r'(\d+)月(\d+)日[—\-](\d+)月(\d+)日', period)
    if m:
        return f'{m.group(1)}.{m.group(2)}-{m.group(3)}.{m.group(4)}'
    return period


def update_crm_date(shape, period_str, crm_date_full):
    _ym = re.match(r'(\d{4})', crm_date_full)
    _year = _ym.group(1) if _ym else '2026'

    for para in shape.text_frame.paragraphs:
        full_text = para.text
        if not full_text:
            continue
        new_text = full_text
        new_text = re.sub(r'\d+月\d+日[—\-]\d+月\d+日', period_str, new_text)
        new_text = re.sub(r'20xx\.[xX\d]+\.[xX\d]+[\-—]20xx\.[xX\d]+\.[xX\d]+', crm_date_full, new_text)
        new_text = re.sub(r'\d{4}\.[xX\d]+\.[xX\d]+[\-—]\d{4}\.[xX\d]+\.[xX\d]+', crm_date_full, new_text)
        new_text = re.sub(r'\d{4}\.\d+\.\d+-\d{4}\.\d+\.\d+', crm_date_full, new_text)
        new_text = new_text.replace('20xx', _year)
        if new_text != full_text and para.runs:
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run.text = ''


# ========== 图表生成（调用 scripts/gen_biz_chart.py） ==========

SCRIPTS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'scripts')
GEN_BIZ_CHART = os.path.join(SCRIPTS_DIR, 'gen_biz_chart.py')


def _generate_biz_chart(eff_data, output_path, title_prefix, date_str, year='2026'):
    chart_items = []
    for item in eff_data:
        name = str(item[0] or '').strip()
        amt = float(item[1] or 0)
        tgt = float(item[2] or 0)
        cnt = 0
        if len(item) > 3 and item[3] is not None:
            try:
                cnt = int(float(item[3]))
            except (ValueError, TypeError):
                cnt = 0
        if not name or tgt <= 0:
            continue
        if name in ('合计', '总计', '小计', '行业合计', '商业合计'):
            continue
        chart_items.append({
            'name': name,
            'amount': round(amt, 2),
            'target': round(tgt, 2),
            'count': cnt
        })

    if not chart_items:
        return False

    json_path = output_path + '.json'
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(chart_items, f, ensure_ascii=False, indent=2)

    try:
        import subprocess
        cmd = [
            sys.executable, GEN_BIZ_CHART,
            '--data', json_path,
            '--output', output_path,
            '--date', date_str,
            '--title-prefix', title_prefix,
            '--title-suffix', '有效商机纳管情况',
            '--subtitle', '有效商机 / 商机储备目标',
            '--xlabel', '有效商机完成率',
            '--target-line', '30',
            '--year', year,
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=30, encoding='utf-8')
        if result.returncode != 0:
            print(f'Chart gen error: {result.stderr}')
            return False
        return os.path.exists(output_path)
    except Exception as e:
        print(f'Chart generation failed: {e}')
        return False
    finally:
        try:
            os.unlink(json_path)
        except Exception:
            pass


# ========== 数据预览（供 routes 调用） ==========

def preview_data_regions(data_bytes, data_map=None):
    """预览 PPT 模块从数据 Excel 中读取的各区域内容，用于校准数据匹配。

    Args:
        data_bytes: 数据 Excel 文件字节
        data_map: 可选的数据映射覆盖

    Returns:
        dict: 预览数据，含各区域行列数和样本
    """
    tmp = tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False)
    tmp.write(data_bytes)
    tmp.close()

    try:
        wb = load_workbook(tmp.name, data_only=True)
        ws = wb.active
    except Exception as e:
        os.unlink(tmp.name)
        raise Exception(f'读取失败: {str(e)}')

    try:
        dm = dict(DEFAULT_DATA_MAP)
        if data_map:
            dm.update(data_map)

        date_c, date_r = _parse_cell_ref(dm['date_cell'])
        period_c, period_r = _parse_cell_ref(dm['period_cell'])
        b27_c, b27_r = _parse_cell_ref(dm['B27_cell'])
        b28_c, b28_r = _parse_cell_ref(dm['B28_cell'])
        j27_c, j27_r = _parse_cell_ref(dm['J27_cell'])
        j28_c, j28_r = _parse_cell_ref(dm['J28_cell'])
        ai27_c, ai27_r = _parse_cell_ref(dm['AI27_cell'])
        ai28_c, ai28_r = _parse_cell_ref(dm['AI28_cell'])

        preview = {
            'date': str(cell_val(ws, date_c, date_r)),
            'period': str(cell_val(ws, period_c, period_r)),
            'B27': str(cell_val(ws, b27_c, b27_r))[:100],
            'B28': str(cell_val(ws, b28_c, b28_r))[:100],
            'J27': str(cell_val(ws, j27_c, j27_r))[:100],
            'J28': str(cell_val(ws, j28_c, j28_r))[:100],
            'AI27': str(cell_val(ws, ai27_c, ai27_r))[:100],
            'AI28': str(cell_val(ws, ai28_c, ai28_r))[:100],
        }

        ir = _parse_range_ref(dm['industry_reserve'])
        cr = _parse_range_ref(dm['commercial_reserve'])
        ie = _parse_range_ref(dm['industry_effective'])
        ce = _parse_range_ref(dm['commercial_effective'])
        ip = _parse_range_ref(dm['industry_progress'])
        cp = _parse_range_ref(dm['commercial_progress'])
        idr = _parse_range_ref(dm['industry_delivered'])
        cdr = _parse_range_ref(dm['commercial_delivered'])

        ind_reserve = read_range(ws, ir[0], ir[1], ir[2], ir[3])
        comm_reserve = read_range(ws, cr[0], cr[1], cr[2], cr[3])
        ind_effective = read_range(ws, ie[0], ie[1], ie[2], ie[3])
        comm_effective = read_range(ws, ce[0], ce[1], ce[2], ce[3])
        ind_progress = read_range(ws, ip[0], ip[1], ip[2], ip[3])
        comm_progress = read_range(ws, cp[0], cp[1], cp[2], cp[3])
        ind_delivered = read_range(ws, idr[0], idr[1], idr[2], idr[3])
        comm_delivered = read_range(ws, cdr[0], cdr[1], cdr[2], cdr[3])

        def to_str_rows(rows):
            return [[str(c) if c is not None else '' for c in r] for r in rows]

        def sample_rows(rows, max_n=3):
            return to_str_rows(rows[:max_n])

        preview['ranges'] = {
            'industry_reserve': {'rows': len(ind_reserve), 'sample': sample_rows(ind_reserve)},
            'commercial_reserve': {'rows': len(comm_reserve), 'sample': sample_rows(comm_reserve)},
            'industry_effective': {'rows': len(ind_effective), 'sample': sample_rows(ind_effective), 'full': to_str_rows(ind_effective)},
            'commercial_effective': {'rows': len(comm_effective), 'sample': sample_rows(comm_effective), 'full': to_str_rows(comm_effective)},
            'industry_progress': {'rows': len(ind_progress), 'sample': sample_rows(ind_progress)},
            'commercial_progress': {'rows': len(comm_progress), 'sample': sample_rows(comm_progress)},
            'industry_delivered': {'rows': len(ind_delivered), 'sample': sample_rows(ind_delivered)},
            'commercial_delivered': {'rows': len(comm_delivered), 'sample': sample_rows(comm_delivered)},
        }

        wb.close()
        os.unlink(tmp.name)
        return preview
    except Exception:
        wb.close()
        os.unlink(tmp.name)
        raise


# ========== 核心：PPT 生成 ==========

def generate_ppt(template_bytes, data_bytes, custom_texts=None, data_map=None):
    """根据模板 PPT 和数据 Excel 生成通报 PPT。"""
    dm = dict(DEFAULT_DATA_MAP)
    if data_map:
        dm.update(data_map)

    date_col, date_row = _parse_cell_ref(dm['date_cell'])
    period_col, period_row = _parse_cell_ref(dm['period_cell'])
    crm_col, crm_row = _parse_cell_ref(dm['crm_date_cell'])
    b27_col, b27_row = _parse_cell_ref(dm['B27_cell'])
    b28_col, b28_row = _parse_cell_ref(dm['B28_cell'])
    j27_col, j27_row = _parse_cell_ref(dm['J27_cell'])
    j28_col, j28_row = _parse_cell_ref(dm['J28_cell'])
    ai27_col, ai27_row = _parse_cell_ref(dm['AI27_cell'])
    ai28_col, ai28_row = _parse_cell_ref(dm['AI28_cell'])

    tmp_ppt = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    tmp_ppt.write(template_bytes)
    tmp_ppt.close()

    tmp_data = tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False)
    tmp_data.write(data_bytes)
    tmp_data.close()

    tmp_out = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    tmp_out.close()

    try:
        wb = load_workbook(tmp_data.name, data_only=True)
        ws = wb.active
    except Exception as e:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except OSError: pass
        return {'ok': False, 'error': f'读取数据文件失败: {str(e)}'}

    try:
        DATE_STR = excel_date_to_str(cell_val(ws, date_col, date_row))
        PERIOD_STR = str(cell_val(ws, period_col, period_row))
        PERIOD_CRM = period_to_crm(PERIOD_STR)

        _year_match = re.match(r'(\d{4})年', DATE_STR)
        YEAR_FULL = _year_match.group(1) if _year_match else '2026'
        YEAR_SHORT = YEAR_FULL[2:]

        CRM_DATE_C30 = str(cell_val(ws, crm_col, crm_row)).strip()
        if CRM_DATE_C30 and CRM_DATE_C30 != 'None':
            CRM_DATE_FULL = CRM_DATE_C30
        else:
            PERIOD_CRM = period_to_crm(PERIOD_STR)
            CRM_DATE_FULL = f'{YEAR_FULL}.{PERIOD_CRM}'.replace('-', f'-{YEAR_FULL}.')

        B27_TEXT = str(cell_val(ws, b27_col, b27_row)).strip().replace('\n', '')
        B28_TEXT = str(cell_val(ws, b28_col, b28_row)).strip().replace('\n', '')
        J27_TEXT = str(cell_val(ws, j27_col, j27_row)).strip()
        J28_TEXT = str(cell_val(ws, j28_col, j28_row)).strip()
        AI27_TEXT = str(cell_val(ws, ai27_col, ai27_row)).strip()
        AI28_TEXT = str(cell_val(ws, ai28_col, ai28_row)).strip()

        if custom_texts:
            B27_TEXT = custom_texts.get('B27', B27_TEXT)
            B28_TEXT = custom_texts.get('B28', B28_TEXT)
            J27_TEXT = custom_texts.get('J27', J27_TEXT)
            J28_TEXT = custom_texts.get('J28', J28_TEXT)
            AI27_TEXT = custom_texts.get('AI27', AI27_TEXT)
            AI28_TEXT = custom_texts.get('AI28', AI28_TEXT)

        ir_r1, ir_r2, ir_c1, ir_c2 = _parse_range_ref(dm['industry_reserve'])
        industry_reserve = read_range(ws, ir_r1, ir_r2, ir_c1, ir_c2)
        industry_reserve.sort(key=lambda x: float(x[6]) if x[6] is not None else 0, reverse=True)
        cr_r1, cr_r2, cr_c1, cr_c2 = _parse_range_ref(dm['commercial_reserve'])
        commercial_reserve = read_range(ws, cr_r1, cr_r2, cr_c1, cr_c2)
        commercial_reserve.sort(key=lambda x: float(x[6]) if x[6] is not None else 0, reverse=True)

        ie_r1, ie_r2, ie_c1, ie_c2 = _parse_range_ref(dm['industry_effective'])
        industry_effective = read_range(ws, ie_r1, ie_r2, ie_c1, ie_c2)
        ce_r1, ce_r2, ce_c1, ce_c2 = _parse_range_ref(dm['commercial_effective'])
        commercial_effective = read_range(ws, ce_r1, ce_r2, ce_c1, ce_c2)
        ip_r1, ip_r2, ip_c1, ip_c2 = _parse_range_ref(dm['industry_progress'])
        industry_progress = read_range(ws, ip_r1, ip_r2, ip_c1, ip_c2)
        cp_r1, cp_r2, cp_c1, cp_c2 = _parse_range_ref(dm['commercial_progress'])
        commercial_progress = read_range(ws, cp_r1, cp_r2, cp_c1, cp_c2)
        id_r1, id_r2, id_c1, id_c2 = _parse_range_ref(dm['industry_delivered'])
        industry_delivered = read_range(ws, id_r1, id_r2, id_c1, id_c2)
        cd_r1, cd_r2, cd_c1, cd_c2 = _parse_range_ref(dm['commercial_delivered'])
        commercial_delivered = read_range(ws, cd_r1, cd_r2, cd_c1, cd_c2)

        wb.close()
    except Exception as e:
        wb.close()
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except OSError: pass
        return {'ok': False, 'error': f'解析数据失败: {str(e)}'}

    # 生成图表
    chart_dir = os.path.join(PPT_DATA_DIR, 'charts')
    os.makedirs(chart_dir, exist_ok=True)
    industry_chart_path = os.path.join(chart_dir, 'industry_chart.png')
    commercial_chart_path = os.path.join(chart_dir, 'commercial_chart.png')

    for old_chart in [industry_chart_path, commercial_chart_path]:
        if os.path.exists(old_chart):
            try: os.unlink(old_chart)
            except OSError: pass

    date_suffix = PERIOD_STR.split('—')[-1] if '—' in PERIOD_STR else PERIOD_STR.split('-')[-1] if '-' in PERIOD_STR else PERIOD_STR

    ind_chart_ok = _generate_biz_chart(industry_effective, industry_chart_path, '行业各分局', date_suffix, year=YEAR_FULL)
    comm_chart_ok = _generate_biz_chart(commercial_effective, commercial_chart_path, '商业各分局', date_suffix, year=YEAR_FULL)

    def _eff_data_debug(eff_data, max_n=5):
        items = []
        for item in eff_data[:max_n]:
            name = str(item[0] or '').strip()
            amt = float(item[1] or 0)
            tgt = float(item[2] or 0)
            items.append(f'{name}: {amt}/{tgt}')
        return '; '.join(items) + (f' ...共{len(eff_data)}行' if len(eff_data) > max_n else '')

    chart_debug = {
        'industry_effective_range': dm['industry_effective'],
        'industry_data_preview': _eff_data_debug(industry_effective),
        'industry_chart_ok': ind_chart_ok,
        'commercial_effective_range': dm['commercial_effective'],
        'commercial_data_preview': _eff_data_debug(commercial_effective),
        'commercial_chart_ok': comm_chart_ok,
    }

    try:
        prs = Presentation(tmp_ppt.name)
    except Exception as e:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except OSError: pass
        return {'ok': False, 'error': f'读取PPT模板失败: {str(e)}'}

    # 模板结构验证
    if len(prs.slides) < 15:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except OSError: pass
        return {'ok': False, 'error': f'模板页面不足（{len(prs.slides)}页，需15页）。请使用正确的商机通报模板。'}

    required_shapes = {
        0: ['date'],
        2: ['表格 14'],
        5: ['Text 7', 'Text 49'],
        6: ['Text 7', 'Text 49'],
    }
    missing = []
    for slide_idx, names in required_shapes.items():
        existing = {s.name for s in prs.slides[slide_idx].shapes}
        for name in names:
            if name not in existing:
                missing.append(f'Slide{slide_idx + 1}/{name}')
    if missing:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except OSError: pass
        return {'ok': False, 'error': f'模板缺少关键元素：{", ".join(missing)}。请使用正确的商机通报模板(含表格、图表占位等)。'}

    _fill_all_slides(prs, DATE_STR, PERIOD_STR, CRM_DATE_FULL, YEAR_FULL, YEAR_SHORT,
                     B27_TEXT, B28_TEXT, J27_TEXT, J28_TEXT, AI27_TEXT, AI28_TEXT,
                     industry_reserve, commercial_reserve,
                     industry_effective, commercial_effective,
                     industry_progress, commercial_progress,
                     industry_delivered, commercial_delivered,
                     industry_chart_path, commercial_chart_path)

    try:
        prs.save(tmp_out.name)
    except Exception as e:
        for f in [tmp_ppt.name, tmp_data.name, tmp_out.name]:
            try: os.unlink(f)
            except OSError: pass
        return {'ok': False, 'error': f'保存PPT失败: {str(e)}'}

    try: os.unlink(tmp_ppt.name)
    except OSError: pass
    try: os.unlink(tmp_data.name)
    except OSError: pass

    ind_below30_count, _ = count_below30(industry_effective)
    comm_below30_count, _ = count_below30(commercial_effective)

    return {
        'ok': True,
        'output_path': tmp_out.name,
        'date': DATE_STR,
        'period': PERIOD_STR,
        'industry_bureaus': len(industry_reserve),
        'commercial_bureaus': len(commercial_reserve),
        'industry_below30': ind_below30_count,
        'commercial_below30': comm_below30_count,
        'chart_debug': chart_debug,
    }


def _fill_all_slides(prs, DATE_STR, PERIOD_STR, CRM_DATE_FULL, YEAR_FULL, YEAR_SHORT,
                      B27_TEXT, B28_TEXT, J27_TEXT, J28_TEXT, AI27_TEXT, AI28_TEXT,
                      industry_reserve, commercial_reserve,
                      industry_effective, commercial_effective,
                      industry_progress, commercial_progress,
                      industry_delivered, commercial_delivered,
                      industry_chart_path, commercial_chart_path):
    """填充 PPT 所有 slide 内容。"""
    # Slide 1: 日期
    for shape in prs.slides[0].shapes:
        if shape.name == 'date':
            set_shape_single_text(shape, DATE_STR, font_name='微软雅黑', font_size=18, font_bold=True)

    _fill_reserve_slide(prs.slides[2], f'1、行业板块---{YEAR_FULL}年商机储备情况（{PERIOD_STR}）',
                        industry_reserve, B27_TEXT, PERIOD_STR, CRM_DATE_FULL)
    _fill_reserve_slide(prs.slides[3], f'1、商校板块---{YEAR_FULL}年商机储备情况（{PERIOD_STR}）',
                        commercial_reserve, B28_TEXT, PERIOD_STR, CRM_DATE_FULL)

    _fill_effective_slide(prs.slides[5], f'2、行业板块---有效商机纳管情况（{YEAR_SHORT}年{PERIOD_STR}）',
                          industry_effective, PERIOD_STR, CRM_DATE_FULL)
    if os.path.exists(industry_chart_path):
        replace_picture(prs.slides[5], '图片 6', industry_chart_path)

    _fill_effective_slide(prs.slides[6], f'2、商业板块---有效商机纳管情况（{YEAR_SHORT}年{PERIOD_STR}）',
                          commercial_effective, PERIOD_STR, CRM_DATE_FULL, is_commercial=True)
    if os.path.exists(commercial_chart_path):
        replace_picture(prs.slides[7 - 1], '图片 1', commercial_chart_path)

    _fill_progress_slide(prs.slides[9], f'3、行业板块--商机项目推进情况（{YEAR_SHORT}年{PERIOD_STR}）',
                         industry_progress, is_industry=True)
    _fill_progress_slide(prs.slides[10], f'3、商校板块--商机项目推进情况（{YEAR_SHORT}年{PERIOD_STR}）',
                         commercial_progress, is_industry=False)

    _fill_delivered_slide(prs.slides[12], f'3、行业板块--{YEAR_SHORT}年已完成交付项目情况（{YEAR_SHORT}年{PERIOD_STR}）',
                          industry_delivered, AI27_TEXT)
    _fill_delivered_slide(prs.slides[13], f'3、商校板块--{YEAR_SHORT}年已完成交付项目情况（{YEAR_SHORT}年{PERIOD_STR}）',
                          commercial_delivered, AI28_TEXT)

    # 全量扫描替换 CRM 日期
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = para.text
                if txt and ('商机创建时间' in txt or 'CRM取数口径' in txt or '20xx' in txt or re.search(r'\d{4}\.[xX\d]+\.[xX\d]+', txt)):
                    update_crm_date(shape, PERIOD_STR, CRM_DATE_FULL)
                    break


def _fill_reserve_slide(slide, title, reserve_data, b_text, PERIOD_STR, CRM_DATE_FULL):
    """填充储备情况 slide（Slide 3/4 通用）。"""
    for shape in slide.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, title, **TITLE_FONT)
        elif shape.name == '表格 14':
            tbl = shape.table
            for i, rd in enumerate(reserve_data):
                r = i + 1
                set_cell_text(tbl.cell(r, 0), str(i + 1))
                for c in range(1, 7):
                    v = rd[c]
                    if c == 6:
                        set_cell_text(tbl.cell(r, c), fmt_pct(v))
                    elif c == 5:
                        set_cell_text(tbl.cell(r, c), fmt_num(v))
                    elif c in (3, 4):
                        set_cell_text(tbl.cell(r, c), fmt_num(v, 0))
                    else:
                        set_cell_text(tbl.cell(r, c), str(v) if v else '')
        elif shape.name == '文本框 4':
            b3 = get_bottom3(reserve_data)
            set_shape_multiline(shape, [
                f'1、{b_text}',
                f'2、储备完成率不足的分局：{b3}'
            ], font_color_line1=None, font_color_line2=(0x00, 0x80, 0x00))
        elif shape.name == '文本框 5':
            update_crm_date(shape, PERIOD_STR, CRM_DATE_FULL)


def _fill_effective_slide(slide, title, eff_data, PERIOD_STR, CRM_DATE_FULL, is_commercial=False):
    """填充有效商机 slide（Slide 6/7 通用）。"""
    below30_count, severe = count_below30(eff_data)
    sum_amt, sum_tgt, rate_pct, rate_text, bureau_count = calc_effective_summary(eff_data)
    names_str = _format_severe_names(severe)

    for shape in slide.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, title, **TITLE_FONT)
        elif shape.name == 'Text 7':
            set_shape_single_text(shape, fmt_num(sum_amt) + '万')
        elif shape.name == 'Text 9':
            set_shape_multiline(shape, ['有效商机/商机储备目标', rate_text])
        elif shape.name == 'Text 49':
            if is_commercial:
                line1 = f'{below30_count}个分局均未达30%储备率（共{bureau_count}个）'
                line2 = f'其中{names_str}商机纳管储备严重不足'
            else:
                line1 = f'{below30_count}个分局未达30%储备率（共{bureau_count}个）'
                line2 = f'其中{names_str}商机储备纳管严重不足'
            set_shape_multiline(shape, [line1, line2], font_size=14)
        elif shape.name in ('文本框 123', '文本框 5'):
            update_crm_date(shape, PERIOD_STR, CRM_DATE_FULL)


def _fill_progress_slide(slide, title, progress_data, is_industry=True):
    """填充项目推进 slide（Slide 10/11 通用）。"""
    s = calc_progress_summary(progress_data)
    for shape in slide.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, title, **TITLE_FONT)
        elif shape.name == '表格 27' or shape.name == '表格 1':
            tbl = shape.table
            for i, rd in enumerate(progress_data):
                r = i + 2
                for c in range(11):
                    val = rd[c] if rd[c] is not None else ''
                    if 6 <= c <= 10:
                        set_cell_text(tbl.cell(r, c), fmt_ca(val))
                    else:
                        set_cell_text(tbl.cell(r, c), str(val))
        elif shape.name in ('文本框 3', '文本框 5'):
            block_name = '行业板块' if is_industry else '商校板块'
            set_shape_multiline(shape, [
                f'{block_name}在途有效商机{s["total"]}条，其中处于售前{s["sq"]}条、售中{s["sz"]}条。',
                f'{block_name}在途有效商机类型分布情况：IDC、互联网专线商机数量为0',
                f'ICT商机数量为{s["ict_c"]}，金额{fmt_comma(s["ict_a"])}万元，基础业务数量为{s["base_c"]}，金额{fmt_comma(s["base_a"])}万元，产数资源商机数量为{s["res_c"]}，金额{fmt_comma(s["res_a"])}万元。'
            ])


def _fill_delivered_slide(slide, title, delivered_data, ai_text):
    """填充交付情况 slide（Slide 13/14 通用）。"""
    count_colors = _delivered_color_indices(delivered_data, 2)
    amount_colors = _delivered_color_indices(delivered_data, 3)

    for shape in slide.shapes:
        if shape.name == '标题 1':
            set_shape_single_text(shape, title, **TITLE_FONT)
        elif shape.name in ('表格 27', '表格 2'):
            tbl = shape.table
            for i, rd in enumerate(delivered_data):
                r = i + 1
                for c in range(4):
                    v = rd[c]
                    color = None
                    if c == 2:
                        color = count_colors.get(i)
                    elif c == 3:
                        color = amount_colors.get(i)
                    fs = 14 if color else 11
                    if c == 3:
                        set_cell_text(tbl.cell(r, c), fmt_num(v), font_color=color, font_size=fs)
                    elif c == 2:
                        set_cell_text(tbl.cell(r, c), fmt_num(v, 0), font_color=color, font_size=fs)
                    else:
                        set_cell_text(tbl.cell(r, c), str(v) if v else '', font_color=color, font_size=fs)
        elif shape.name in ('文本框 3', '文本框 5'):
            set_shape_single_text(shape, ai_text)
